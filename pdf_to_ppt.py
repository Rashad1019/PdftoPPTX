import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io
import re  # For extracting numbers from filenames

# Function to extract numerical part from filename for sorting
def extract_number(filename):
    match = re.search(r'\d+', filename)
    return int(match.group()) if match else float('inf')  # If no number, put at end

# Function to convert multiple PDFs to a single PPTX
def convert_pdfs_to_pptx(pdf_files):
    # Sort the PDF files based on the numerical part in their names
    sorted_pdfs = sorted(pdf_files, key=lambda f: extract_number(f.name))

    # Create a new PowerPoint presentation
    prs = Presentation()

    # Set slide dimensions to match a standard widescreen (16:9) aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for pdf_file in sorted_pdfs:
        # Read PDF bytes
        pdf_bytes = pdf_file.read()

        # Open the PDF from bytes
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")

        for page_num in range(len(pdf_document)):
            # Render PDF page as image
            page = pdf_document.load_page(page_num)
            pix = page.get_pixmap(dpi=300)  # High DPI for better quality
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            # Save image to bytes
            img_bytes = io.BytesIO()
            img.save(img_bytes, format="PNG")
            img_bytes.seek(0)

            # Add a blank slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Add the image to fill the slide while maintaining aspect ratio
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            img_width, img_height = img.size

            # Calculate scaling factor to fill the slide completely
            scale_x = slide_width / img_width
            scale_y = slide_height / img_height
            scale = max(scale_x, scale_y)  # Use the larger scale to fill the slide

            # Calculate new dimensions (will be larger than slide in one dimension)
            new_width = img_width * scale
            new_height = img_height * scale

            # Center the image on the slide (parts may be cropped)
            left = (slide_width - new_width) / 2
            top = (slide_height - new_height) / 2

            slide.shapes.add_picture(img_bytes, left, top, width=new_width, height=new_height)

        pdf_document.close()  # Clean up

    # Save PPTX to bytes
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)

    return pptx_bytes

# Streamlit app
st.title("PDF to PowerPoint Converter")

uploaded_files = st.file_uploader("Upload PDF files (e.g., PDF0.pdf to PDF10.pdf)", type=["pdf"], accept_multiple_files=True)

if uploaded_files:
    st.write(f"Uploaded {len(uploaded_files)} PDF files.")

    if st.button("Convert to PPTX"):
        if len(uploaded_files) == 0:
            st.error("Please upload at least one PDF file.")
        else:
            with st.spinner("Converting..."):
                pptx_bytes = convert_pdfs_to_pptx(uploaded_files)

            # Provide download button
            st.download_button(
                label="Download PPTX",
                data=pptx_bytes,
                file_name="converted.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            st.success("Conversion complete! Download your PPTX file.")
else:
    st.info("Please upload PDF files to begin.")